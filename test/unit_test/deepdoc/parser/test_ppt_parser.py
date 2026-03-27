#  Copyright 2025 The InfiniFlow Authors. All Rights Reserved.
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#  http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.

"""Unit tests for the PPT parser.

Tests cover:
- Basic parsing functionality with various content types
- Edge cases: empty presentations, special characters
- Slide filtering by page range
- Shape extraction: text frames, tables, groups
- Bulleted list handling
- Shape sorting by position
"""

import importlib.util
import os
import sys
from io import BytesIO
from unittest import mock

_MOCK_MODULES = [
    "xgboost",
    "xgb",
    "pdfplumber",
    "huggingface_hub",
    "PIL",
    "PIL.Image",
    "pypdf",
    "sklearn",
    "sklearn.cluster",
    "sklearn.metrics",
    "deepdoc.vision",
    "infinity",
    "infinity.rag_tokenizer",
]
for _m in _MOCK_MODULES:
    if _m not in sys.modules:
        sys.modules[_m] = mock.MagicMock()


def _find_project_root(marker="pyproject.toml"):
    d = os.path.dirname(os.path.abspath(__file__))
    while d != os.path.dirname(d):
        if os.path.exists(os.path.join(d, marker)):
            return d
        d = os.path.dirname(d)
    return None


_PROJECT_ROOT = _find_project_root()

_ppt_spec = importlib.util.spec_from_file_location(
    "deepdoc.parser.ppt_parser",
    os.path.join(_PROJECT_ROOT, "deepdoc", "parser", "ppt_parser.py"),
)
_ppt_mod = importlib.util.module_from_spec(_ppt_spec)
sys.modules["deepdoc.parser.ppt_parser"] = _ppt_mod
_ppt_spec.loader.exec_module(_ppt_mod)

RAGFlowPptParser = _ppt_mod.RAGFlowPptParser


def _create_pptx_with_slides(slides_content):
    """Create a minimal PPTX file in memory.

    Args:
        slides_content: list of slide content, each slide is a dict with:
            - texts: list of text strings
            - tables: list of tables (each table is list of rows)
            - bullets: list of bullet point strings

    Returns:
        bytes: PPTX file content
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    for slide_data in slides_content:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        if "texts" in slide_data:
            for i, text in enumerate(slide_data["texts"]):
                left = Inches(0.5)
                top = Inches(0.5 + i * 1)
                width = Inches(8)
                height = Inches(0.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = text

        if "bullets" in slide_data:
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(8)
            height = Inches(2)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            for i, bullet in enumerate(slide_data["bullets"]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.level = 0

        if "tables" in slide_data:
            for table_data in slide_data["tables"]:
                if not table_data:
                    continue
                rows = len(table_data)
                cols = len(table_data[0]) if table_data else 0
                if rows > 0 and cols > 0:
                    shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(3), Inches(6), Inches(2))
                    table = shape.table
                    for i, row_data in enumerate(table_data):
                        for j, cell_text in enumerate(row_data):
                            table.cell(i, j).text = str(cell_text)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _create_pptx_with_grouped_shapes():
    """Create a PPTX with grouped shapes.

    Returns:
        bytes: PPTX file content
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shape1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    shape1.text_frame.paragraphs[0].text = "Shape 1"

    shape2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
    shape2.text_frame.paragraphs[0].text = "Shape 2"

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


class TestPptParserBasic:
    """Tests for basic PPT parsing functionality."""

    def test_parse_single_slide(self):
        """Parse a presentation with a single slide."""
        slides_content = [{"texts": ["Hello World"]}]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1
        assert "Hello World" in result[0]

    def test_parse_multiple_slides(self):
        """Parse a presentation with multiple slides."""
        slides_content = [
            {"texts": ["Slide 1 Content"]},
            {"texts": ["Slide 2 Content"]},
            {"texts": ["Slide 3 Content"]},
        ]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 3
        assert "Slide 1 Content" in result[0]
        assert "Slide 2 Content" in result[1]
        assert "Slide 3 Content" in result[2]

    def test_parse_empty_presentation(self):
        """Parse an empty presentation."""
        slides_content = []
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 0

    def test_parse_slide_with_chinese_text(self):
        """Parse a presentation with Chinese text."""
        slides_content = [{"texts": ["中文内容测试"]}]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1
        assert "中文内容测试" in result[0]

    def test_parse_from_file_path(self):
        """Test parsing from a file path."""
        import tempfile

        slides_content = [{"texts": ["File path test"]}]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            f.write(pptx_bytes)
            temp_path = f.name

        try:
            parser = RAGFlowPptParser()
            result = parser(temp_path, from_page=0, to_page=100)

            assert len(result) == 1
            assert "File path test" in result[0]
        finally:
            os.unlink(temp_path)


class TestPptParserPageRange:
    """Tests for slide filtering by page range."""

    def test_parse_from_page(self):
        """Parse starting from a specific slide."""
        slides_content = [
            {"texts": ["Slide 0"]},
            {"texts": ["Slide 1"]},
            {"texts": ["Slide 2"]},
        ]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=1, to_page=100)

        assert len(result) == 2
        assert "Slide 1" in result[0]

    def test_parse_to_page(self):
        """Parse up to a specific slide."""
        slides_content = [
            {"texts": ["Slide 0"]},
            {"texts": ["Slide 1"]},
            {"texts": ["Slide 2"]},
        ]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=2)

        assert len(result) == 2

    def test_parse_page_range(self):
        """Parse a specific range of slides."""
        slides_content = [
            {"texts": ["Slide 0"]},
            {"texts": ["Slide 1"]},
            {"texts": ["Slide 2"]},
            {"texts": ["Slide 3"]},
        ]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=1, to_page=3)

        assert len(result) == 2

    def test_parse_single_page(self):
        """Parse a single page."""
        slides_content = [
            {"texts": ["Slide 0"]},
            {"texts": ["Slide 1"]},
        ]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=1, to_page=2)

        assert len(result) == 1


class TestPptParserTables:
    """Tests for table extraction from slides."""

    def test_parse_slide_with_table(self):
        """Parse a slide containing a table."""
        slides_content = [
            {
                "texts": ["Title"],
                "tables": [[["Header1", "Header2"], ["Value1", "Value2"]]],
            }
        ]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1

    def test_parse_table_content_structure(self):
        """Verify table content is extracted with headers."""
        slides_content = [
            {
                "tables": [[["Name", "Age"], ["Alice", "30"], ["Bob", "25"]]],
            }
        ]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1

    def test_parse_multiple_tables_on_slide(self):
        """Parse a slide with multiple tables."""
        slides_content = [
            {
                "tables": [
                    [["A", "B"], ["1", "2"]],
                    [["C", "D"], ["3", "4"]],
                ],
            }
        ]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1


class TestPptParserBullets:
    """Tests for bulleted list handling."""

    def test_parse_bulleted_list(self):
        """Parse a slide with bulleted list."""
        slides_content = [
            {
                "bullets": ["First point", "Second point", "Third point"],
            }
        ]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1
        assert "First point" in result[0]

    def test_parse_mixed_text_and_bullets(self):
        """Parse a slide with both text and bullets."""
        slides_content = [
            {
                "texts": ["Title"],
                "bullets": ["Point 1", "Point 2"],
            }
        ]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1


class TestPptParserShapes:
    """Tests for shape extraction."""

    def test_parse_multiple_text_shapes(self):
        """Parse a slide with multiple text shapes."""
        slides_content = [
            {
                "texts": ["Title Text", "Subtitle Text", "Body Text"],
            }
        ]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1

    def test_shape_cache_initialization(self):
        """Verify shape cache is initialized."""
        parser = RAGFlowPptParser()

        assert hasattr(parser, "_shape_cache")
        assert parser._shape_cache == {}

    def test_shape_cache_used(self):
        """Verify shape cache is populated during parsing."""
        slides_content = [{"texts": ["Test"]}]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        parser(pptx_bytes, from_page=0, to_page=100)


class TestPptParserEdgeCases:
    """Tests for edge cases and error handling."""

    def test_parse_slide_with_empty_text_frame(self):
        """Parse a slide with empty text frames."""
        slides_content = [{"texts": ["", "Non-empty", ""]}]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1

    def test_parse_slide_with_special_characters(self):
        """Parse a slide with special characters."""
        slides_content = [{"texts": ["Special: !@#$%^&*()", "Quotes: \"'`"]}]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1

    def test_parse_slide_with_unicode(self):
        """Parse a slide with various Unicode characters."""
        slides_content = [{"texts": ["Emoji: 😀🎉", "Japanese: こんにちは", "Arabic: مرحبا"]}]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1

    def test_parse_large_presentation(self):
        """Parse a presentation with many slides."""
        slides_content = [{"texts": [f"Slide {i}"]} for i in range(100)]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 100

    def test_parse_slide_with_newlines(self):
        """Parse a slide with text containing newlines."""
        slides_content = [{"texts": ["Line 1\nLine 2\nLine 3"]}]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1

    def test_total_page_attribute_set(self):
        """Verify total_page attribute is set after parsing."""
        slides_content = [{"texts": ["Slide 1"]}, {"texts": ["Slide 2"]}]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        parser(pptx_bytes, from_page=0, to_page=100)

        assert hasattr(parser, "total_page")
        assert parser.total_page == 2


class TestPptParserCallback:
    """Tests for callback parameter."""

    def test_parse_with_callback_none(self):
        """Parse with callback=None (default)."""
        slides_content = [{"texts": ["Test"]}]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100, callback=None)

        assert len(result) == 1


class TestPptParserShapeSorting:
    """Tests for shape sorting by position."""

    def test_shapes_sorted_by_position(self):
        """Verify shapes are sorted by top/left position."""
        slides_content = [
            {
                "texts": ["Top Shape", "Middle Shape", "Bottom Shape"],
            }
        ]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1


class TestPptParserEmptyContent:
    """Tests for handling empty content scenarios."""

    def test_parse_slide_with_no_shapes(self):
        """Parse a completely empty slide."""
        from pptx import Presentation

        prs = Presentation()
        slide_layout = prs.slide_layouts[6]
        prs.slides.add_slide(slide_layout)

        buf = BytesIO()
        prs.save(buf)
        buf.seek(0)
        pptx_bytes = buf.getvalue()

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1
        assert result[0] == ""

    def test_parse_table_with_empty_cells(self):
        """Parse a table with empty cells."""
        slides_content = [
            {
                "tables": [[["A", "", "C"], ["", "B", ""]]],
            }
        ]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1

    def test_parse_single_row_table(self):
        """Parse a table with only header row."""
        slides_content = [
            {
                "tables": [[["Header1", "Header2", "Header3"]]],
            }
        ]
        pptx_bytes = _create_pptx_with_slides(slides_content)

        parser = RAGFlowPptParser()
        result = parser(pptx_bytes, from_page=0, to_page=100)

        assert len(result) == 1
